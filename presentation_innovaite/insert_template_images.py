#!/usr/bin/env python3
"""Insert content images from the Beamer deck into the InnovAIte PPT template."""

from pptx import Presentation
from pptx.util import Inches

INPUT = "InovAIte_ppt_template_SVK_v2.pptx"
OUTPUT = "InovAIte_ppt_template_with_images.pptx"

prs = Presentation(INPUT)

# Slide numbers are 1-based in comments, 0-based in python-pptx.
# Placement is chosen to preserve the template title/text/footer areas.
placements = [
    # slide 2: Research goal / problem context
    (1, "assets/smellhier.png", 0.95, 3.65, 4.95, None),
    # slide 3: Dependency model
    (2, "assets/depgraph.png", 1.10, 2.45, None, 3.95),
    # slide 4: Code processing pipeline
    (3, "assets/pipeline.png", 0.95, 5.50, 11.20, None),
]

for slide_idx, image, left, top, width, height in placements:
    slide = prs.slides[slide_idx]
    kwargs = {
        "left": Inches(left),
        "top": Inches(top),
    }
    if width is not None:
        kwargs["width"] = Inches(width)
    if height is not None:
        kwargs["height"] = Inches(height)
    slide.shapes.add_picture(image, **kwargs)

prs.save(OUTPUT)
print(f"wrote {OUTPUT}")
